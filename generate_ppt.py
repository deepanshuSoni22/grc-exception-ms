"""
GRC Exception & Policy Waiver Management System
PowerPoint Generator for Societe Generale Hackathon 2026
Team: ByteBuilders
Members: Vinutha S T, Serah Joyce Rasquinha

Run this script to generate the full presentation automatically.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# CONFIGURATION
# ============================================================
OUTPUT_FILE = "ByteBuilders_GRC_Presentation.pptx"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Scheme - Corporate Dark Blue + Gold
COLORS = {
    "primary": RGBColor(26, 42, 58),      # #1a2a3a - Dark Blue
    "secondary": RGBColor(212, 168, 67),  # #d4a843 - Gold
    "accent": RGBColor(45, 65, 85),        # #2d4155 - Lighter Blue
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
    "red": RGBColor(200, 50, 50),
    "green": RGBColor(50, 180, 80),
    "yellow": RGBColor(220, 180, 50),
    "gray": RGBColor(150, 150, 150),
    "light_gray": RGBColor(240, 240, 240),
}

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def add_title_slide(prs, title, subtitle, team_names, event):
    """Create a professional title slide with branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background - Dark Blue
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["primary"]
    bg.line.fill.background()
    
    # Gold Accent Bar - Top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["secondary"]
    bar.line.fill.background()
    
    # Gold Accent Bar - Bottom
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_HEIGHT - Inches(0.15), SLIDE_WIDTH, Inches(0.15)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = COLORS["secondary"]
    bar2.line.fill.background()
    
    # Team Name (Top Right)
    team_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(3.5), Inches(0.4), Inches(3), Inches(0.6)
    )
    tf = team_box.text_frame
    tf.text = "BYTEBUILDERS"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["secondary"]
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Main Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(11), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["white"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.8), Inches(11), Inches(1)
    )
    tf = sub_box.text_frame
    tf.text = subtitle
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = COLORS["white"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Team Members
    members_box = slide.shapes.add_textbox(
        Inches(1), Inches(5), Inches(11), Inches(0.8)
    )
    tf = members_box.text_frame
    tf.text = f"Team Members: {team_names}"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Event
    event_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.8), Inches(11), Inches(0.6)
    )
    tf = event_box.text_frame
    tf.text = event
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = COLORS["secondary"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Decorative circle - bottom right
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, SLIDE_WIDTH - Inches(2.5), SLIDE_HEIGHT - Inches(2.5), 
        Inches(1.5), Inches(1.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS["secondary"]
    circle.line.fill.background()
    circle.fill.transparency = 0.6


def add_content_slide(prs, title, content_items, bg_color=None, 
                      title_color=None, accent_bar=True):
    """Create a content slide with bullet points and optional accent bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg_color = bg_color or COLORS["white"]
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    
    # Accent Bar (Left side)
    if accent_bar:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SLIDE_HEIGHT
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["secondary"]
        bar.line.fill.background()
    
    # Title
    title_color = title_color or COLORS["primary"]
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(11), Inches(1)
    )
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = title_color
    
    # Gold underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(2.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["secondary"]
    line.line.fill.background()
    
    # Content
    y_pos = 1.8
    for item in content_items:
        if isinstance(item, dict):
            # Rich content with formatting
            text = item.get("text", "")
            bold = item.get("bold", False)
            color = item.get("color", COLORS["black"])
            level = item.get("level", 0)
            font_size = item.get("size", 20)
        else:
            text = item
            bold = False
            color = COLORS["black"]
            level = 0
            font_size = 20
        
        # Bullet or text box
        if text.startswith("• "):
            text = text[2:]
            box = slide.shapes.add_textbox(
                Inches(0.8 + (level * 0.4)), Inches(y_pos), 
                Inches(11.5 - (level * 0.4)), Inches(0.6)
            )
            tf = box.text_frame
            tf.text = "• " + text
            tf.paragraphs[0].font.size = Pt(font_size)
            tf.paragraphs[0].font.color.rgb = color
            if bold:
                tf.paragraphs[0].font.bold = True
            y_pos += 0.5
        else:
            box = slide.shapes.add_textbox(
                Inches(0.8 + (level * 0.4)), Inches(y_pos), 
                Inches(11.5 - (level * 0.4)), Inches(0.6)
            )
            tf = box.text_frame
            tf.text = text
            tf.paragraphs[0].font.size = Pt(font_size)
            tf.paragraphs[0].font.color.rgb = color
            if bold:
                tf.paragraphs[0].font.bold = True
            y_pos += 0.5
        
        # Check if we need a new page
        if y_pos > 6.5:
            break


def add_slide_with_two_columns(prs, title, left_content, right_content):
    """Create a slide with two columns of content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["white"]
    bg.line.fill.background()
    
    # Accent Bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SLIDE_HEIGHT
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["secondary"]
    bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(11), Inches(1)
    )
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["primary"]
    
    # Left Column
    left_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8), Inches(5.5), Inches(5)
    )
    tf = left_box.text_frame
    tf.text = left_content
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = COLORS["black"]
    
    # Right Column
    right_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.8), Inches(5.5), Inches(5)
    )
    tf = right_box.text_frame
    tf.text = right_content
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = COLORS["black"]


def add_architecture_slide(prs):
    """Create the architecture diagram slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 247, 250)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.3), Inches(11), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.text = "SOLUTION ARCHITECTURE — THE MVP"
    tf.paragraphs[0].font.size = Pt(34)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["primary"]
    
    # Gold underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(3), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["secondary"]
    line.line.fill.background()
    
    # LAYER 1: Presentation Layer
    layer1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), 
        Inches(12), Inches(1.5)
    )
    layer1.fill.solid()
    layer1.fill.fore_color.rgb = COLORS["primary"]
    layer1.line.fill.background()
    
    layer1_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.6), Inches(11.6), Inches(1.3)
    )
    tf = layer1_text.text_frame
    tf.text = "PRESENTATION LAYER\nDashboard • Registry • Alerts • Audit Report"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["white"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.color.rgb = RGBColor(200, 200, 200)
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Arrow
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(6), Inches(3.1), Inches(0.8), Inches(0.5)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS["secondary"]
    arrow.line.fill.background()
    
    # LAYER 2: Application Layer
    layer2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.7), 
        Inches(12), Inches(1.5)
    )
    layer2.fill.solid()
    layer2.fill.fore_color.rgb = COLORS["accent"]
    layer2.line.fill.background()
    
    layer2_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(3.8), Inches(11.6), Inches(1.3)
    )
    tf = layer2_text.text_frame
    tf.text = "APPLICATION LAYER\nFlask Backend • Risk Engine • Scheduler • Anomaly Detection"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["white"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.color.rgb = RGBColor(200, 200, 200)
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Arrow
    arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(6), Inches(5.3), Inches(0.8), Inches(0.5)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLORS["secondary"]
    arrow2.line.fill.background()
    
    # LAYER 3: Data Layer
    layer3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.9), 
        Inches(12), Inches(1.2)
    )
    layer3.fill.solid()
    layer3.fill.fore_color.rgb = COLORS["primary"]
    layer3.line.fill.background()
    
    layer3_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(6.0), Inches(11.6), Inches(1.0)
    )
    tf = layer3_text.text_frame
    tf.text = "DATA LAYER\nSQLite Database — 600 Records (Apr 2025–Apr 2026)"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["white"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.color.rgb = RGBColor(200, 200, 200)
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER


def add_risk_scoring_slide(prs):
    """Create the risk scoring engine slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["white"]
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.3), Inches(11), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.text = "RISK SCORING ENGINE — QUANTIFYING RISK"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["primary"]
    
    # Left Column - Factors
    left_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.3), Inches(6), Inches(5.5)
    )
    tf = left_box.text_frame
    tf.text = """MULTI-FACTOR RISK SCORE (0-100)

┌─────────────────────────────────────┐
│ Exception Type          │ 40 pts   │
│ Declared Risk Level     │ 45 pts   │
│ Duration (>365 days)    │ 20 pts   │
│ Zombie Status           │ 20 pts   │
│ Justification Quality   │  5 pts   │
└─────────────────────────────────────┘

THRESHOLDS:
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
  LOW      MEDIUM    HIGH    CRITICAL
  0───44───│───65───│───85───│──100
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━"""
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLORS["black"]
    
    # Right Column - Anomalies
    right_box = slide.shapes.add_textbox(
        Inches(7), Inches(1.3), Inches(5.5), Inches(5.5)
    )
    tf = right_box.text_frame
    tf.text = """ANOMALY DETECTION — 5 RULES

┌─────────────────────────────────────┐
│ EXPIRED_ACTIVE        │ 🔴 CRITICAL│
│ CRITICAL_RISK         │ 🔴 CRITICAL│
│ HIGH_RISK_LONG        │ 🟠 HIGH    │
│ LONG_RUNNING          │ 🟠 HIGH    │
│ STALLED_REVIEW        │ 🟡 MEDIUM  │
└─────────────────────────────────────┘

EACH ANOMALY TRIGGERS:
• Alert in Monitor
• Automated Email
• Action Recommendation
• Audit Trail Entry"""
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLORS["black"]


def add_closing_slide(prs):
    """Create a professional closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background - Dark Blue
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["primary"]
    bg.line.fill.background()
    
    # Gold Accent Bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["secondary"]
    bar.line.fill.background()
    
    # Main Text
    main_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(11), Inches(1.5)
    )
    tf = main_box.text_frame
    tf.text = "THANK YOU"
    tf.paragraphs[0].font.size = Pt(60)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["white"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Team
    team_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(11), Inches(0.8)
    )
    tf = team_box.text_frame
    tf.text = "BYTEBUILDERS"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["secondary"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Members
    members_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.3), Inches(11), Inches(0.6)
    )
    tf = members_box.text_frame
    tf.text = "Vinutha S T  |  Serah Joyce Rasquinha"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.3), Inches(11), Inches(0.8)
    )
    tf = tagline_box.text_frame
    tf.text = '"We don\'t just track exceptions — we eliminate the risk they create."'
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = COLORS["secondary"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


# ============================================================
# MAIN GENERATION FUNCTION
# ============================================================

def generate_presentation():
    """Generate the complete 17-slide presentation."""
    
    print("🚀 Generating presentation for ByteBuilders...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # ===== SLIDE 1: TITLE =====
    print("  Generating Slide 1: Title...")
    add_title_slide(
        prs,
        "GRC Process Exception & Policy Waiver Management System",
        "An Enterprise-Grade Governance, Risk & Compliance Solution",
        "Vinutha S T  |  Serah Joyce Rasquinha",
        "Societe Generale Hackathon 2026"
    )
    
    # ===== SLIDE 2: CIA TRIAD =====
    print("  Generating Slide 2: CIA Triad...")
    add_content_slide(
        prs,
        "THE CIA TRIAD — THE PROBLEM FRAMEWORK",
        [
            {"text": "CONFIDENTIALITY 🔒", "bold": True, "color": COLORS["primary"], "size": 22},
            {"text": "• 'Admin access granted for 3 days — active for 3 years'", "level": 1, "size": 18},
            {"text": "• 30% of breaches exploit policy exceptions", "level": 1, "size": 18},
            {"text": "", "size": 10},
            {"text": "INTEGRITY 📋", "bold": True, "color": COLORS["primary"], "size": 22},
            {"text": "• 'Temporary firewall rule — 18 months later, still open'", "level": 1, "size": 18},
            {"text": "• Policy deviations become undocumented reality", "level": 1, "size": 18},
            {"text": "", "size": 10},
            {"text": "AVAILABILITY ⚡", "bold": True, "color": COLORS["primary"], "size": 22},
            {"text": "• New hire can't access systems — exception granted, never reviewed", "level": 1, "size": 18},
            {"text": "• Emergency exceptions bypass governance", "level": 1, "size": 18},
            {"text": "", "size": 10},
            {"text": "➡ When exceptions aren't managed, ALL THREE pillars collapse.", "bold": True, "color": COLORS["red"], "size": 20},
        ]
    )
    
    # ===== SLIDE 3: BUSINESS PROBLEM =====
    print("  Generating Slide 3: Business Problem...")
    add_content_slide(
        prs,
        "THE BUSINESS PROBLEM — ENTERPRISE REALITY",
        [
            {"text": "The Uncomfortable Truth:", "bold": True, "color": COLORS["primary"], "size": 22},
            {"text": "100+ policies ✅ — Reality deviates from policy 🔴 constantly", "size": 18},
            {"text": "", "size": 10},
            {"text": "CHAOS IN NUMBERS:", "bold": True, "color": COLORS["primary"], "size": 20},
            {"text": "• Exceptions tracked: 📧 Emails, Slack threads, Excel sheets", "level": 1, "size": 17},
            {"text": "• 'Temporary' exceptions: 5 years+ active", "level": 1, "size": 17},
            {"text": "• Visibility: ❌ Nonexistent — no one knows how many exist", "level": 1, "size": 17},
            {"text": "• Audit confidence: ❌ 'Which exceptions are still justified?'", "level": 1, "size": 17},
            {"text": "", "size": 10},
            {"text": "REAL FAILURES:", "bold": True, "color": COLORS["red"], "size": 20},
            {"text": "🔴 Case 1: Exception granted 'temporarily' → 3 years later → breach", "level": 1, "size": 17},
            {"text": "🔴 Case 2: 15 exceptions pending review for '3-6 months' → forgotten", "level": 1, "size": 17},
            {"text": "🔴 Case 3: 5 vendors with undocumented exceptions → compliance nightmare", "level": 1, "size": 17},
        ]
    )
    
    # ===== SLIDE 4: STRIDE =====
    print("  Generating Slide 4: STRIDE...")
    add_slide_with_two_columns(
        prs,
        "STRIDE THREAT MODEL — OUR COUNTERMEASURES",
        """S - Spoofing
• Centralized authenticated registry
• Session-based login + SHA-256 hashing

T - Tampering
• Audit trail + immutable log
• notification_log.txt + database

R - Repudiation
• Non-repudiation proof
• Every action logged with timestamp & user

I - Information Disclosure
• Visibility into ALL exceptions
• Dashboard + real-time alerts""",
        """D - Denial of Service
• Automated expiry enforcement
• Scheduler + revocation alerts

E - Elevation of Privilege
• Risk scoring + anomaly detection
• 5-rule engine catches admin abuse

━━━━━━━━━━━━━━━━━━━━━━━━━━━
➡ We don't just track exceptions —
   we neutralize every STRIDE threat.
━━━━━━━━━━━━━━━━━━━━━━━━━━━"""
    )
    
    # ===== SLIDE 5: Framework Alignment =====
    print("  Generating Slide 5: Framework Alignment...")
    add_content_slide(
        prs,
        "FRAMEWORK ALIGNMENT — NIST, GDPR, CIS",
        [
            {"text": "NIST SP 800-53", "bold": True, "color": COLORS["primary"], "size": 22},
            {"text": "• AC-2 (Account Management): Exceptions cannot circumvent controls", "level": 1, "size": 17},
            {"text": "• PL-4 (Rules of Behavior): Justification required; short = higher risk", "level": 1, "size": 17},
            {"text": "• AU-2 (Audit Events): Complete audit trail: who, what, when, why", "level": 1, "size": 17},
            {"text": "", "size": 10},
            {"text": "GDPR Article 25 — Data Protection by Design", "bold": True, "color": COLORS["primary"], "size": 22},
            {"text": "✅ Exceptions are exceptional, tracked, and justified", "level": 1, "size": 17},
            {"text": "✅ Expired exceptions flagged immediately", "level": 1, "size": 17},
            {"text": "✅ Automated revocation alerts", "level": 1, "size": 17},
            {"text": "", "size": 10},
            {"text": "CIS Controls 1.1 — Inventory of IT Assets", "bold": True, "color": COLORS["primary"], "size": 22},
            {"text": "✅ 100% of exceptions centralized, searchable, audit-exportable", "level": 1, "size": 17},
            {"text": "✅ Exceptions treated as assets — authorized/unauthorized tracked", "level": 1, "size": 17},
        ]
    )
    
    # ===== SLIDE 6: Architecture =====
    print("  Generating Slide 6: Architecture...")
    add_architecture_slide(prs)
    
    # ===== SLIDE 7: Risk Scoring =====
    print("  Generating Slide 7: Risk Scoring...")
    add_risk_scoring_slide(prs)
    
    # ===== SLIDE 8: Dashboard =====
    print("  Generating Slide 8: Dashboard...")
    add_content_slide(
        prs,
        "DASHBOARD — REAL-TIME VISIBILITY",
        [
            {"text": "8 KPI CARDS — AT A GLANCE", "bold": True, "color": COLORS["primary"], "size": 20},
            {"text": "Total: 600  |  Active: 182  |  Expiring: 12  |  Expired-Not Revoked: 3", "size": 17},
            {"text": "CRITICAL: 23  |  HIGH: 78  |  PENDING: 15  |  REVOKED: 45", "size": 17},
            {"text": "", "size": 10},
            {"text": "VISUAL ANALYTICS:", "bold": True, "color": COLORS["primary"], "size": 20},
            {"text": "• Risk Level Donut Chart — CRITICAL / HIGH / MEDIUM / LOW", "level": 1, "size": 17},
            {"text": "• Exception Type Bar Chart — Admin / Firewall / Encryption / Data Access", "level": 1, "size": 17},
            {"text": "• Top Requesters by Accumulated Risk — Surfaces hidden exposure", "level": 1, "size": 17},
            {"text": "", "size": 10},
            {"text": "TOP RISK HOTSPOTS:", "bold": True, "color": COLORS["red"], "size": 20},
            {"text": "1. USR-1234 — Admin access to Production (5 months, no review)", "level": 1, "size": 17},
            {"text": "2. Partner API — GL System access (90 days, no renewal)", "level": 1, "size": 17},
            {"text": "3. Legacy System — Encryption waiver (2+ years, unsustainable)", "level": 1, "size": 17},
        ]
    )
    
    # ===== SLIDE 9: Alerts =====
    print("  Generating Slide 9: Alerts...")
    add_content_slide(
        prs,
        "ALERT MONITOR — PROACTIVE NOTIFICATIONS",
        [
            {"text": "5 ALERT CATEGORIES — AUTOMATICALLY DETECTED", "bold": True, "color": COLORS["primary"], "size": 20},
            {"text": "🔴 Expired — Not Revoked: 3  →  REVOKE IMMEDIATELY", "level": 1, "size": 17},
            {"text": "🟡 Expiring Soon: 12  →  RENEW OR REVOKE", "level": 1, "size": 17},
            {"text": "🟠 Long Running: 45  →  SCHEDULE RENEWAL REVIEW", "level": 1, "size": 17},
            {"text": "🟠 Stalled Review: 8  →  ESCALATE TO APPROVER", "level": 1, "size": 17},
            {"text": "🔴 Critical Risk: 23  →  REVIEW WITH CISO", "level": 1, "size": 17},
            {"text": "", "size": 10},
            {"text": "AUTOMATED NOTIFICATIONS — ZERO MANUAL TRIGGERS", "bold": True, "color": COLORS["primary"], "size": 20},
            {"text": "• On Flask startup — fires within 5 seconds", "level": 1, "size": 17},
            {"text": "• Daily at 08:00 — production schedule", "level": 1, "size": 17},
            {"text": "", "size": 10},
            {"text": "WHAT GETS SENT:", "bold": True, "color": COLORS["primary"], "size": 18},
            {"text": "📧 Requester Summary  |  📧 Approver Summary  |  📧 Admin Digest", "level": 1, "size": 17},
            {"text": "✅ Deduplication built-in — no duplicate emails on restart", "level": 1, "size": 17},
        ]
    )
    
    # ===== SLIDE 10: Audit =====
    print("  Generating Slide 10: Audit...")
    add_content_slide(
        prs,
        "AUDIT READINESS — ONE-CLICK COMPLIANCE",
        [
            {"text": "EXECUTIVE SUMMARY (Auto-Generated)", "